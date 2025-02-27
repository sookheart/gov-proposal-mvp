import streamlit as st
import anthropic
import json
import time
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
import re

# 제목과 설명 표시
st.title("🏢 정부지원과제 사업계획서 작성기")
st.write(
    "정부지원과제 공고 내용과 회사 정보를 입력하면 Claude AI가 맞춤형 사업계획서를 작성해 드립니다. "
    "이 앱을 사용하려면 Anthropic의 Claude API 키가 필요합니다."
)

def split_into_sections(text):
    """사업계획서 텍스트를 섹션별로 분리합니다."""
    # 주요 섹션 패턴 (숫자. 제목)
    section_pattern = r'\n\s*\d+\.\s+[^\n]+'
    sections = re.split(section_pattern, text)
    
    # 섹션 제목 추출
    section_titles = re.findall(section_pattern, text)
    
    # 첫 번째 섹션은 일반적으로 서문이므로 "서문" 제목 추가
    if sections and not section_titles:
        sections = sections[1:]  # 빈 첫 섹션 제거
    elif sections and len(sections) > len(section_titles):
        sections[0] = sections[0].strip()
        if sections[0]:  # 첫 섹션에 내용이 있으면
            section_titles.insert(0, "서문")
        else:
            sections = sections[1:]  # 빈 첫 섹션 제거
    
    # 섹션 제목과 내용 결합
    result = []
    for i in range(len(section_titles)):
        if i < len(sections):
            section_content = sections[i].strip()
            if section_content:  # 내용이 있는 경우에만 추가
                result.append((section_titles[i].strip(), section_content))
    
    return result

def create_pptx(title, sections):
    """사업계획서 내용으로 PowerPoint 파일을 생성합니다."""
    prs = Presentation()
    
    # 제목 슬라이드 추가
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    
    # 각 섹션을 별도 슬라이드로 추가
    for section_title, section_content in sections:
        # 섹션 제목 슬라이드
        section_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(section_slide_layout)
        slide.shapes.title.text = section_title
        
        # 섹션 내용을 여러 슬라이드로 분할 (텍스트가 많을 경우)
        content_chunks = [section_content[i:i+800] for i in range(0, len(section_content), 800)]
        
        # 첫 번째 청크는 섹션 제목 슬라이드에 추가
        if content_chunks:
            tf = slide.shapes.placeholders[1].text_frame
            tf.text = content_chunks[0]
        
        # 나머지 청크는 새 슬라이드에 추가
        for chunk in content_chunks[1:]:
            content_slide_layout = prs.slide_layouts[2]
            slide = prs.slides.add_slide(content_slide_layout)
            slide.shapes.title.text = f"{section_title} (계속)"
            tf = slide.shapes.placeholders[1].text_frame
            tf.text = chunk
    
    # 메모리에 저장
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    
    return output

# 사용자로부터 Claude API 키 입력 받기
claude_api_key = st.text_input("Claude API Key", type="password")

if not claude_api_key:
    st.info("계속하려면 Claude API 키를 입력해주세요.", icon="🔑")
else:
    try:
        # Claude 클라이언트 생성
        client = anthropic.Anthropic(api_key=claude_api_key)
        
        # 사용 가능한 모델 목록 가져오기
        with st.spinner("사용 가능한 Claude 모델을 확인 중입니다..."):
            models_response = client.models.list()
            available_models = [model.id for model in models_response.data if "claude" in model.id.lower()]
        
        if not available_models:
            st.error("사용 가능한 Claude 모델이 없습니다. API 키를 확인해주세요.")
        else:
            # 사용자가 모델 선택
            selected_model = st.selectbox("사용할 Claude 모델 선택", available_models)
            
            # 입력 섹션 생성
            st.subheader("1. 정부지원과제 공고 정보")
            announcement = st.text_area(
                "정부지원과제 공고 내용을 입력해주세요",
                height=200,
                placeholder="예: 중소기업 디지털 전환 지원사업 공고\n지원대상: 제조업 분야 중소기업\n지원내용: 스마트공장 구축 및 디지털 전환..."
            )
            
            st.subheader("2. 회사 정보")
            company_name = st.text_input("회사명")
            business_type = st.text_input("업종")
            company_size = st.selectbox("회사 규모", ["소기업", "중기업", "중견기업"])
            employee_count = st.number_input("직원 수", min_value=1, value=10)
            annual_revenue = st.number_input("연간 매출액(백만원)", min_value=0, value=500)
            
            company_history = st.text_area(
                "회사 연혁",
                height=100,
                placeholder="예: 2015년 설립, 2018년 벤처기업 인증, 2020년 수출유망기업 선정..."
            )
            
            core_business = st.text_area(
                "주요 사업 내용",
                height=100,
                placeholder="회사의 주요 제품, 서비스, 기술력 등을 설명해주세요."
            )
            
            company_strength = st.text_area(
                "회사의 강점",
                height=100,
                placeholder="기술력, 인력, 특허, 시장 점유율 등 회사의 강점을 설명해주세요."
            )
            
            project_goals = st.text_area(
                "이 정부지원과제를 통해 달성하고자 하는 목표",
                height=100,
                placeholder="예: 생산성 향상, 신제품 개발, 해외 시장 진출 등"
            )
            
            # 사업계획서 생성 버튼
            if st.button("사업계획서 생성"):
                if not announcement or not company_name or not business_type or not core_business:
                    st.error("필수 정보(공고 내용, 회사명, 업종, 주요 사업 내용)를 모두 입력해주세요.")
                else:
                    # 회사 정보 구성
                    company_info = f"""
                    회사명: {company_name}
                    업종: {business_type}
                    회사 규모: {company_size}
                    직원 수: {employee_count}명
                    연간 매출액: {annual_revenue}백만원
                    회사 연혁: {company_history}
                    주요 사업 내용: {core_business}
                    회사의 강점: {company_strength}
                    프로젝트 목표: {project_goals}
                    """
                    
                    # 프롬프트 작성
                    prompt = f"""
                    당신은 정부지원과제 사업계획서 작성 전문가입니다. 다음 정보를 기반으로 약 30페이지 분량의 상세한 사업계획서를 작성해주세요.
                    
                    ## 정부지원과제 공고 내용:
                    {announcement}
                    
                    ## 회사 정보:
                    {company_info}
                    
                    사업계획서는 다음 구조를 따라 작성해주세요:
                    
                    1. 사업 개요
                       - 사업의 배경 및 필요성
                       - 사업의 목표 및 범위
                       - 기대효과
                    
                    2. 기업 소개
                       - 기업 현황 및 연혁
                       - 조직 구성 및 인력 현황
                       - 주요 사업 분야 및 실적
                       - 기술력 및 경쟁력 분석
                    
                    3. 시장 분석
                       - 국내외 시장 동향
                       - 경쟁사 분석
                       - SWOT 분석
                    
                    4. 사업 내용
                       - 핵심 기술 소개
                       - 추진 전략 및 방법론
                       - 세부 추진 계획 및 일정
                       - 추진 체계
                    
                    5. 사업화 전략
                       - 마케팅 전략
                       - 판매 전략
                       - 해외 진출 전략 (해당 시)
                    
                    6. 소요 자금 계획
                       - 총 사업비 구성
                       - 연차별 자금 소요 계획
                       - 정부지원금 활용 계획
                    
                    7. 기대효과 및 성과지표
                       - 기술적 기대효과
                       - 경제적 기대효과
                       - 성과 측정 지표 및 방법
                    
                    8. 향후 발전 방향
                       - 중장기 발전 계획
                       - 후속 연구개발 계획
                    
                    각 항목을 상세히 작성하고, 공고 내용과 회사 정보를 적절히 연계하여 설득력 있는 사업계획서를 작성해주세요.
                    문서는 한글로 작성해주세요.
                    
                    각 섹션과 소제목을 명확하게 구분하여 작성해주세요. 각 섹션 시작은 '숫자. 제목' 형식으로 해주세요.
                    """
                    
                    with st.spinner("사업계획서 작성 중... (약 3-5분 소요됩니다)"):
                        try:
                            # 스트리밍 응답을 위한 빈 텍스트 공간 생성
                            response_placeholder = st.empty()
                            full_response = ""
                            
                            # Claude API를 사용하여 사업계획서 생성 (스트리밍 방식)
                            with client.messages.stream(
                                model=selected_model,
                                max_tokens=64000,  # Claude 모델의 최대 출력 토큰 수 제한
                                temperature=0.7,
                                messages=[
                                    {"role": "user", "content": prompt}
                                ]
                            ) as stream:
                                for text in stream.text_stream:
                                    full_response += text
                                    response_placeholder.markdown(full_response)
                            
                            # 결과 표시
                            st.subheader("📝 생성된 사업계획서")
                            
                            # 탭 생성 (텍스트 전체 보기 / 섹션별 보기)
                            tab1, tab2 = st.tabs(["텍스트 전체 보기", "섹션별 보기"])
                            
                            with tab1:
                                st.markdown(full_response)
                                
                                # 다운로드 버튼 추가 (텍스트)
                                st.download_button(
                                    label="사업계획서 다운로드 (.txt)",
                                    data=full_response,
                                    file_name=f"{company_name}_사업계획서.txt",
                                    mime="text/plain"
                                )
                            
                            with tab2:
                                # 섹션별로 분리
                                sections = split_into_sections(full_response)
                                
                                for i, (section_title, section_content) in enumerate(sections):
                                    with st.expander(f"{section_title}", expanded=i==0):
                                        st.markdown(section_content)
                                
                                # 섹션별 파일 생성
                                section_files = {}
                                for section_title, section_content in sections:
                                    clean_title = section_title.strip().replace('.', '_').replace(' ', '_')
                                    section_files[f"{clean_title}.txt"] = section_content
                                
                                # 다운로드 버튼 (섹션별 텍스트)
                                for file_name, content in section_files.items():
                                    st.download_button(
                                        label=f"{file_name} 다운로드",
                                        data=content,
                                        file_name=file_name,
                                        mime="text/plain",
                                        key=file_name
                                    )
                            
                            # PowerPoint 생성 및 다운로드
                            try:
                                ppt_title = f"{company_name} 사업계획서"
                                ppt_buffer = create_pptx(ppt_title, sections)
                                
                                st.subheader("📊 PowerPoint 다운로드")
                                st.download_button(
                                    label="사업계획서 PowerPoint 다운로드 (.pptx)",
                                    data=ppt_buffer,
                                    file_name=f"{company_name}_사업계획서.pptx",
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                                )
                            except Exception as e:
                                st.warning(f"PowerPoint 생성 중 오류가 발생했습니다: {str(e)}")
                                st.info("Python-pptx 라이브러리 설치가 필요합니다. requirements.txt 파일에 'python-pptx'를 추가해주세요.")
                            
                        except Exception as e:
                            st.error(f"사업계획서 생성 중 오류가 발생했습니다: {str(e)}")
                
    except Exception as e:
        st.error(f"API 연결 중 오류가 발생했습니다: {str(e)}")
